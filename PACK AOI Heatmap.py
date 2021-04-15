# - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
#   Program Summary
#   Purpose
#   Inputs : None
#   Outputs : BOT & TOP heatmap, BOT & TOP heatmap with board image, pptx with 
#             BOT & TOP heatmap with board image and defect information inside. (JZ 04142021)
#   Side Effects : Copies contents of Source directory to Destination Directory
# - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -

import datetime
import os
import pandas as pd
import re
import time
import cx_Oracle 

pd.set_option('display.max_columns', None)
pd.set_option('display.width', None)
pd.set_option('display.max_rows', None)

cx_Oracle.init_oracle_client(lib_dir=r"C:\JZhou\oracle sql\instantclient_19_10")#Change the direction to where you put the oracle instant client
#For PANGUS PRD account, ask Xiaolong MA for credential.
dsnStr = cx_Oracle.makedsn("10.20.193.53", "1521", "IMSDBPROD")#The structure is (ip,port,SID)
db = cx_Oracle.connect('xiaoma','foxconn123MKE',dsn=dsnStr)#The structure is (Username,Password,dsnstr)

cur = db.cursor() 

################ Change the week and date only ##################
week = 'WK18'
date = "3-29"
start_date = "'03/01/2021'"
end_date = "'03/31/2021'"

# - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
#           Begin - Function Definitions
# - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
# This function will help to select useful columns from original data.
def read_defect_data(defect_data):

    raw_defect_data = defect_data
    raw_defect_data.columns = "WO code,Product,Label code,Station,Defect name,Location,Repair,Week".split(",")

    raw_defect_data.drop(['WO code', 'Product', 'Station', 'Repair', 'Week'], axis='columns', inplace=True)

    raw_defect_data.rename(columns={"Label code": "Barcode", "Location": "RefDesignator", "Defect name": "Defect",
                                    }, inplace=True)

    defect_data = raw_defect_data[['Barcode', 'RefDesignator', 'Defect']]

    return defect_data

# This function will help to seperate the defects into TOP/BOT two parts based
# on the component data.
def process_data(defect_data, component_data):
    merged_data = pd.merge(defect_data, component_data, on='RefDesignator')
    print(component_data)
    print(defect_data)
    print(merged_data)
    # Check for mis-matched components
    for index, row in defect_data.iterrows():
        if row['RefDesignator'] not in merged_data.values:
            print('\n' + row['RefDesignator'] + ' does not exist in merged_data')

    merged_data.drop(['XOffset', 'YOffset', 'Angle', 'Customer P/N'], axis='columns', inplace=True)
    heatmap_data = merged_data[
        ['Barcode', 'RefDesignator', 'Defect', 'Part', 'Package',
         'Supplier', 'Supplier P/N', 'Hon Hai P/N', 'RegionID']]

    return heatmap_data



# This function mainly collect defect data without time, the SQL command is from 
# Xiaolong's script and the start_date and end_date will be valued while using the function.
def Data_collection(start_date,end_date):
    command = f"SELECT DISTINCT WO_code, PROD_ID AS Product, LB_ID AS Label_code, WP_ID AS Station,  \
                                 BAD_ITEM_NAME AS Defect_name, BAD_POINT AS Location, ' ' AS Repair, CONCAT('Rampup ', 'WKXX') AS Week\
                    FROM\
                (\
                SELECT MO as WO_code, LB_ID, PROD_ID, WP_ID, IS_PASS, TO_CHAR(t1.SYS_CRT_DATE, 'MM/DD/YYYY hh24:mi:ss') Process_completion_time, BAD_ITEM_NAME, BAD_POINT, t2.FIELD_EX1,\
                             RANK()OVER(PARTITION BY LB_ID ORDER BY TO_DATE(TO_CHAR(t1.SYS_CRT_DATE, 'MM/DD/YYYY hh24:mi:ss'), 'MM/DD/YYYY hh24:mi:ss')) rank_time\
                   FROM IMS.TB_PM_QC_HD t1\
                   JOIN IMS.TB_PM_QC_DT t2 ON t1.QC_ID = t2.QC_ID AND t1.SYS_CRT_DATE = t2.SYS_CRT_DATE\
                   JOIN IMS.TB_BS_BAD_ITEM b ON b.BAD_ITEM_ID = t2.BAD_ITEM_ID\
                 WHERE WP_ID = 'PTH-Packing-AOI'\
                      AND LB_ID LIKE 'W%'\
                      AND t1.SYS_CRT_DATE BETWEEN TO_DATE({start_date}, 'MM/DD/YYYY') AND TO_DATE({end_date}, 'MM/DD/YYYY')\
                      AND b.BAD_ITEM_ID NOT IN ('Dimension', 'OCROCV', 'Coplanarity')\
                      AND LB_ID IN (SELECT DISTINCT LB_ID\
                                                    FROM (SELECT LB_ID,  WP_ID, IS_PASS, WP_CMP_DATE\
                                                                   FROM (\
                                                                   SELECT  LB_ID, WP_ID, IS_PASS, WP_CMP_DATE,  MIN (WP_CMP_DATE) OVER (PARTITION BY LB_ID, WP_ID) PROCESS_COMPLETION\
                                                                                  FROM IMS.TB_PM_MO_LBWP\
                                                                               )\
                                                                 WHERE WP_CMP_DATE = PROCESS_COMPLETION\
                                                                      AND WP_CMP_DATE BETWEEN TO_DATE({start_date}, 'MM/DD/YYYY') AND TO_DATE({end_date}, 'MM/DD/YYYY')\
                                                                      AND WP_ID ='PTH-Packing-AOI')\
                                                    WHERE IS_PASS = 'N')\
                )g\
                 WHERE rank_time = 1"
    cur.execute(command) # Execute SQL command
    rows = cur.fetchall()
    data = pd.DataFrame(rows)
    
    return data

# This function mainly collect defect data with time, the SQL command is from 
# Xiaolong's script and the start_date and end_date will be valued while using the function.
def Data_collection_shift(start_date,end_date):
    command = f"SELECT MO as WO_code, LB_ID, PROD_ID, WP_ID, IS_PASS, TO_CHAR(t1.SYS_CRT_DATE, 'MM/DD/YYYY hh24:mi:ss') Process_completion_time, BAD_ITEM_NAME, BAD_POINT, t2.FIELD_EX1,\
                         RANK()OVER(PARTITION BY LB_ID ORDER BY TO_DATE(TO_CHAR(t1.SYS_CRT_DATE, 'MM/DD/YYYY hh24:mi:ss'), 'MM/DD/YYYY hh24:mi:ss')) rank_time\
               FROM IMS.TB_PM_QC_HD t1\
               JOIN IMS.TB_PM_QC_DT t2 ON t1.QC_ID = t2.QC_ID AND t1.SYS_CRT_DATE = t2.SYS_CRT_DATE\
               JOIN IMS.TB_BS_BAD_ITEM b ON b.BAD_ITEM_ID = t2.BAD_ITEM_ID\
             WHERE WP_ID = 'PTH-Packing-AOI'\
                  AND LB_ID LIKE 'W%'\
                  AND t1.SYS_CRT_DATE BETWEEN TO_DATE({start_date}, 'MM/DD/YYYY') AND TO_DATE({end_date}, 'MM/DD/YYYY')\
                  AND b.BAD_ITEM_ID NOT IN ('Dimension', 'OCROCV', 'Coplanarity')\
                  AND LB_ID IN (SELECT DISTINCT LB_ID\
                                                FROM (SELECT LB_ID,  WP_ID, IS_PASS, WP_CMP_DATE\
                                                               FROM (\
                                                               SELECT  LB_ID, WP_ID, IS_PASS, WP_CMP_DATE,  MIN (WP_CMP_DATE) OVER (PARTITION BY LB_ID, WP_ID) PROCESS_COMPLETION\
                                                                              FROM IMS.TB_PM_MO_LBWP\
                                                                           )\
                                                             WHERE WP_CMP_DATE = PROCESS_COMPLETION\
                                                                  AND WP_CMP_DATE BETWEEN TO_DATE({start_date}, 'MM/DD/YYYY') AND TO_DATE({end_date}, 'MM/DD/YYYY')\
                                                                  AND WP_ID ='PTH-Packing-AOI')\
                                                WHERE IS_PASS = 'N')"
    cur.execute(command) # Execute SQL command
    rows = cur.fetchall()
    data = pd.DataFrame(rows)
    data.columns = ["WO_CODE","LB_ID","PROD_ID","WP_ID","IS_PASS","PROCESS_COMPLETION_TIME","BAD_ITEM_NAME","BAD_POINT","FIELD_EX1","RANK_TIME"]
    return data


# - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
#           Begin - Main Program
# - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
#Change route
output_path = r'C:/JZhou/L6/Heatmap'
os.chdir(output_path)

# - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
# Read originial defect data
raw_defect_data = Data_collection(start_date,end_date)
# raw_defect_data = pd.read_excel(os.path.join(output_path,f'Pangus Raw Data/{week}/PACK AOI_{week}_{date}.xlsx'))
# Save the raw data in Pangus raw data folder
os.makedirs(output_path +  f'/Pangus Raw Data/{week}',exist_ok=True)
raw_defect_data.to_excel(os.path.join(output_path,f'Pangus Raw Data/{week}/PACK AOI_{week}_{date}.xlsx'))
defect_data = read_defect_data(raw_defect_data)

###################### for the TOP AOI ##############################################
# Define the station
station = "BOT_AOI" 

# Read the component list
file403 = 'ComponentData_403BOT.csv'
#file403 = 'ComponentData_403TOP.csv'
#file503 = 'ComponentData_503BOT.csv'
#file503 = 'ComponentData_503TOP.csv'

component_data_bot = pd.read_csv(file403)
 
# write BOT heatmap file
heatmap_data_BOT = process_data(defect_data, component_data_bot)

#heatmap_data.to_csv(output_path + '\\' + f'{station}_Region_data' + '.csv', index=False)

###################### for the TOP AOI ##############################################
station = 'TOP_AOI'
file403 = 'ComponentData_403TOP.csv'
component_data_top = pd.read_csv(file403)
heatmap_data_TOP = process_data(defect_data, component_data_top)

########################## Make an Excel Spreadsheet  ####################
os.makedirs(output_path +  f'/Heatmap Raw Data/{week}',exist_ok=True)

# Create a Pandas Excel writer using XlsxWriter as the engine.
writer = pd.ExcelWriter(output_path + f'/Heatmap Raw Data/{week}/{week}_{date}_PACK_Heatmap Data' + '.xlsx', engine='xlsxwriter')

# Write each dataframe to a different worksheet.
heatmap_data_BOT.to_excel(writer, sheet_name='Bot')
heatmap_data_TOP.to_excel(writer, sheet_name='Top')

# Close the Pandas Excel writer and output the Excel file.
writer.save()
time.sleep(0.5)



# - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
###########################     Draw Heatmap     ###############################
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd 
import os
import numpy as np 
sns.set()

###############################################################################
station_list = ['Bot','Top']

for station in station_list:
    # set up the working directory 
    os.chdir(os.path.join(output_path,f'Heatmap Raw Data/{week}' ))
    excel_name = f'{week}_{date}_PACK_Heatmap Data.xlsx'
    
    # data processing: raw data will be stored as top and bot separately in a excel.
    raw_data = pd.read_excel(f'./{excel_name}', sheet_name = f'{station}')
    raw_data1 = raw_data.sort_values(by = ['RegionID','RefDesignator'],ascending = (True,True))
    # raw_data.head()
    # pick up the counts of each regionID and sum them up into a column ,  this number is the m.
    region_counts = raw_data.groupby(['RegionID']).size().reset_index(name='counts')
    
    # sum up all count, this number is the n.
    counts_sum_up = raw_data.groupby(['RegionID']).size().sum()
    #####################################################################
    # make a function to do the list combination 
    def lists_combination(lists, code=','):
        try:
            import reduce
        except:
            from functools import reduce
            
        def myfunc(list1, list2):
            return [str(i)+code+str(j) for i in list1 for j in list2]
        return reduce(myfunc, lists)
    
    region1  = [0,1,2,3,4,5,6,7]
    region2  = [0,1,2,3,4,5,6,7]
    
    # make the regions location coordinates 
    regions = pd.Series(lists_combination([region1,region2], code=','))
    #######################################################################
    
    # make up an empty df to save the occurence data into
    region_defect = pd.DataFrame()
    region_defect["RegionID"] = regions
    region_defect['counts'] = np.nan
    
    # merge the region_defect with the region_counts  
    region_defect = region_defect.merge(region_counts, how='left', on='RegionID')
    region_defect.drop(columns = ['counts_x'],inplace = True)
    region_defect.fillna(0,inplace=True)
    region_defect.rename(columns = {'counts_y':'counts'}, inplace= True)
    
    # make the heatmap data, iloc[ ] re-assign the counts 
    col = '0 1 2 3 4 5 6 7'.split(' ')
    row_index = '0 1 2 3 4 5 6 7'.split(' ')
    
    plt_data =pd.DataFrame(index = pd.Index(row_index), columns=(col)) 
    
    # start to assign the value
    for index,values in enumerate(region_defect['RegionID'].tolist()):
        row_loc = int(values.split(',')[0])  
        col_loc = int(values.split(',')[1]) 
        plt_data.iloc[row_loc,col_loc] = region_defect['counts'].iloc[index]
    
    # reverse the row and fit the heatmap
    plt_data = plt_data.iloc[:,::-1]
    
    # change the data from object into int
    plt_data = plt_data.astype(int)
    plt_data = plt_data.T
    
    # Generate a df which is the same structure with the heatmap. The content
    # inside will be the top three error location in this region and their 
    # occurance time.
    map_information = pd.DataFrame(index = pd.Index(row_index), columns=(col))
    for x in range(plt_data.shape[0]):
        for y in range(plt_data.shape[1]):
            # Change column and row number into regionID
            X = y
            Y = 7 - x
            if f'{X},{Y}' in list(raw_data1['RegionID']):
                # Select data based on regionID
                raw_data_temp = raw_data1[raw_data1['RegionID'] == f'{X},{Y}']
                # Sort RefDesignator based on their occurance.
                duplicate_id = raw_data_temp['RefDesignator'].value_counts()
                duplicate_id = duplicate_id.reset_index()
                duplicate_id.columns = ["RefDesignator","Occur number"]
                # Design the layout for each region
                if duplicate_id.shape[0] > 3:
                    map_information_temp = ''
                    for i in range(3):
                        x1 = duplicate_id["RefDesignator"][i]
                        # If there is _ in the name, then change a new line
                        x2 = list(x1)
                        if '_' in x2:
                            x2 = ["_\n" if i =="_" else i for i in x2]
                            x1 = ''.join(x2)
                        y1 = duplicate_id["Occur number"][i]
                        if i == 2:
                            map_information_temp += f"{x1} {y1}"
                        else:
                            map_information_temp +=f"{x1} {y1}\n"
                else:
                    map_information_temp = ''
                    for i in range(duplicate_id.shape[0]):
                        x1 = duplicate_id["RefDesignator"][i]
                        x2 = list(x1)
                        if '_' in x2:
                            x2 = ["_\n" if i =="_" else i for i in x2]
                            x1 = ''.join(x2)

                        y1 = duplicate_id["Occur number"][i]
                        if i == duplicate_id.shape[0] - 1:
                            map_information_temp += f"{x1} {y1}"
                        else:
                            map_information_temp += f"{x1} {y1}\n"
                            
                map_information.iloc[x,y] = map_information_temp
                    
            else:
                map_information.iloc[x,y] = ''


    # plot the graph
    fig, ax = plt.subplots(figsize=(10, 10))
    
    # * 100 to emphasize the pattern
    sns.heatmap(plt_data*100,annot = map_information,alpha = 0.7 ,fmt="", annot_kws={ 'color':'black'},cmap='OrRd',ax=ax,cbar=False, xticklabels=False,yticklabels=False,
                linewidth=0.7, linecolor='red', square=True)
    
    # label_y = ax.get_yticklabels()
    # plt.setp(label_y, rotation=360, horizontalalignment='right')
    # label_x = ax.get_xticklabels()
    # plt.setp(label_x, horizontalalignment='right')
    plt.show()
    
    
    # save the fig
    os.chdir(os.path.join(output_path,f'Heatmap Raw Data/{week}'))
    fig.savefig(f'{station}_Heatmap_first_layer.png', dpi=600, bbox_inches='tight')
    
    ############################################
    
    # for top only
    if station == 'Top':
        from PIL import Image
        from PIL import ImageOps
        
        img1 = Image.open( f"../../{station}.JPG")
        img1 = img1.convert('RGBA')
         
        img2 = Image.open( f"./{station}_Heatmap_first_layer.png")

        img2 = img2.convert('RGBA')
        
        # Resize foreground down from 500x500 to 100x100
        img2 = img2.resize((824,697))
        
        img = Image.blend(img1, img2, 0.8)
        img.show()
       
        img.save( f"./{station}_{date}_Heatmap.png")
    
    #####################################################
    
    # for bot only
    if station == 'Bot' :
        from PIL import Image
        from PIL import ImageOps
        
        img1 = Image.open( f"../../{station}.JPG")
        img1 = img1.convert('RGBA')
         
        img2 = Image.open( f"./{station}_Heatmap_first_layer.png")

        img2 = img2.convert('RGBA')
        
        # Resize foreground down from 500x500 to 100x100
        img2 = img2.resize((832,707))
        
        img = Image.blend(img1, img2, 0.8)
        img.show()
       
        img.save( f"./{station}_{date}_Heatmap.png")
        
########################## Shift & Occurence Table ############################
time.sleep(1)

Top_TOP3 = pd.read_excel(f'./{excel_name}', sheet_name = 'Top')
Bot_TOP3 = pd.read_excel(f'./{excel_name}', sheet_name = 'Bot')

# Read defect data with time/shift information and save it
sht_data = Data_collection_shift(start_date,end_date)
sht_data_name = os.path.join(output_path,f'Pangus Raw Data/{week}/PACK AOI_{week}_{date}(Shift).xlsx')
sht_data.to_excel(sht_data_name,index = False)
time.sleep(1)

# Generate defect file with time 
shift_data = pd.read_excel(sht_data_name,usecols = ['LB_ID','PROCESS_COMPLETION_TIME','BAD_ITEM_NAME','BAD_POINT'])                                    
shift_data.columns = ['Label code','Process completion time','Defect name','Location']

merged_df1 = Top_TOP3.merge(shift_data,how='inner',left_on='Barcode',right_on='Label code')
merged_df2 = Bot_TOP3.merge(shift_data,how='inner',left_on='Barcode',right_on='Label code')
merged = pd.concat([merged_df1,merged_df2], axis=0,ignore_index=True)

merged['Process completion time'] = pd.to_datetime(merged['Process completion time'])
merged['hour'] = merged['Process completion time'].dt.hour
merged['Date'] = merged['Process completion time'].dt.date
merged['Shift'] = np.where(merged['hour'].isin([0,1,2,3,4,5,6]), '3rd', \
                             np.where(merged['hour'].isin([7,8,9,10,11,12,13,14]),'1st', '2nd'))

# save the file
merged_name = os.path.join(output_path,f'Heatmap Raw Data/{week}/{week}_{date}_Heatmap Date(deep granularity).xlsx')
merged.to_excel(merged_name,index=False)



useful_column = ['Label code','Defect','RefDesignator','RegionID','Process completion time','Shift']
merged_data = merged[useful_column]

##############################################################################
# If replace, start from here.
##############################################################################
merged_data_top = pd.merge(merged_data, component_data_top, on='RefDesignator')
# Check for mis-matched components
for index, row in merged_data.iterrows():
    if row['RefDesignator'] not in merged_data_top.values:
        print('\n' + row['RefDesignator'] + ' does not exist in merged_data')

merged_data_top = merged_data_top[['Label code','Defect','RefDesignator','RegionID_x','Process completion time','Shift']]
merged_data_top.columns = ['Label code','Defect','RefDesignator','Location','Process completion time','Shift']
merged_data_top = merged_data_top.drop_duplicates()


merged_data_bot = pd.merge(merged_data, component_data_bot, on='RefDesignator')
# Check for mis-matched components
for index, row in merged_data.iterrows():
    if row['RefDesignator'] not in merged_data_bot.values:
        print('\n' + row['RefDesignator'] + ' does not exist in merged_data')

merged_data_bot = merged_data_bot[['Label code','Defect','RefDesignator','RegionID_x','Process completion time','Shift']]
merged_data_bot.columns = ['Label code','Defect','RefDesignator','Location','Process completion time','Shift']
merged_data_bot = merged_data_bot.drop_duplicates()

# for i in range(merged_data_bot.shape[0]):
#     if merged_data_bot['RegionID_x'][i] != merged_data_bot['RegionID_y'][i]:
#         print(i)

#  Generate TOP/BOT defect list rank by Defect type and Location occurence.

merged_defect_top = merged_data_top['Defect'].value_counts()
merged_defect_top = merged_defect_top.reset_index()
merged_defect_top.columns = ["Defect name","Occur number"]

if merged_defect_top.shape[0] >= 3:
    defect_list_top = list(merged_defect_top['Defect name'][0:3])
    
    merged_data1_top = merged_data_top[merged_data_top['Defect'] == defect_list_top[0]]
    merged_data1_top = merged_data1_top.reset_index(drop = True)
    merged_data1_top['Occurence'] = 1
    merged_location_top1 = merged_data1_top['RefDesignator'].value_counts()
    merged_location_top1 = merged_location_top1 .reset_index()
    merged_location_top1.columns = ["RefDesignator","Occur number"]
    for i in range(merged_data1_top.shape[0]):
        for j in range(merged_location_top1.shape[0]):
            if merged_data1_top['RefDesignator'][i] == merged_location_top1['RefDesignator'][j]:
                merged_data1_top['Occurence'][i] = merged_location_top1["Occur number"][j]
    merged_data1_top = merged_data1_top.sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
    
    merged_data2_top = merged_data_top[merged_data_top['Defect'] == defect_list_top[1]]
    merged_data2_top = merged_data2_top.reset_index(drop = True)
    merged_data2_top['Occurence'] = 1
    merged_location_top2 = merged_data2_top['RefDesignator'].value_counts()
    merged_location_top2 = merged_location_top2 .reset_index()
    merged_location_top2.columns = ["RefDesignator","Occur number"]
    for i in range(merged_data2_top.shape[0]):
        for j in range(merged_location_top2.shape[0]):
            if merged_data2_top['RefDesignator'][i] == merged_location_top2['RefDesignator'][j]:
                merged_data2_top['Occurence'][i] = merged_location_top2["Occur number"][j]
    merged_data2_top = merged_data2_top.sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
    
    merged_data3_top = merged_data_top[merged_data_top['Defect'] == defect_list_top[2]]
    merged_data3_top = merged_data3_top.reset_index(drop = True)
    merged_data3_top['Occurence'] = 1
    merged_location_top3 = merged_data3_top['RefDesignator'].value_counts()
    merged_location_top3 = merged_location_top3 .reset_index()
    merged_location_top3.columns = ["RefDesignator","Occur number"]
    for i in range(merged_data3_top.shape[0]):
        for j in range(merged_location_top3.shape[0]):
            if merged_data3_top['RefDesignator'][i] == merged_location_top3['RefDesignator'][j]:
                merged_data3_top['Occurence'][i] = merged_location_top3["Occur number"][j]
    merged_data3_top = merged_data3_top.sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
    merged_Data_top = pd.concat([merged_data1_top,merged_data2_top,merged_data3_top])
    merged_Data_top = merged_Data_top.reset_index(drop = True)
elif merged_defect_top.shape[0] == 2:
    defect_list_top = list(merged_defect_top['Defect name'][0:2])
    
    merged_data1_top = merged_data_top[merged_data_top['Defect'] == defect_list_top[0]]
    merged_data1_top = merged_data1_top.reset_index(drop = True)
    merged_data1_top['Occurence'] = 1
    merged_location_top1 = merged_data1_top['RefDesignator'].value_counts()
    merged_location_top1 = merged_location_top1 .reset_index()
    merged_location_top1.columns = ["RefDesignator","Occur number"]
    for i in range(merged_data1_top.shape[0]):
        for j in range(merged_location_top1.shape[0]):
            if merged_data1_top['RefDesignator'][i] == merged_location_top1['RefDesignator'][j]:
                merged_data1_top['Occurence'][i] = merged_location_top1["Occur number"][j]
    merged_data1_top = merged_data1_top.sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
    
    merged_data2_top = merged_data_top[merged_data_top['Defect'] == defect_list_top[1]]
    merged_data2_top = merged_data2_top.reset_index(drop = True)
    merged_data2_top['Occurence'] = 1
    merged_location_top2 = merged_data2_top['RefDesignator'].value_counts()
    merged_location_top2 = merged_location_top2 .reset_index()
    merged_location_top2.columns = ["RefDesignator","Occur number"]
    for i in range(merged_data2_top.shape[0]):
        for j in range(merged_location_top2.shape[0]):
            if merged_data2_top['RefDesignator'][i] == merged_location_top2['RefDesignator'][j]:
                merged_data2_top['Occurence'][i] = merged_location_top2["Occur number"][j]
    merged_data2_top = merged_data2_top.sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
    
    merged_Data_top = pd.concat([merged_data1_top,merged_data2_top])
    merged_Data_top = merged_Data_top.reset_index(drop = True)
elif merged_defect_top.shape[0] == 1:
    defect_list_top = merged_defect_top['Defect name'][0]
    
    merged_data1_top = merged_data_top[merged_data_top['Defect'] == defect_list_top]
    merged_data1_top = merged_data1_top.reset_index(drop = True)
    merged_data1_top['Occurence'] = 1
    merged_location_top1 = merged_data1_top['RefDesignator'].value_counts()
    merged_location_top1 = merged_location_top1 .reset_index()
    merged_location_top1.columns = ["RefDesignator","Occur number"]
    for i in range(merged_data1_top.shape[0]):
        for j in range(merged_location_top1.shape[0]):
            if merged_data1_top['RefDesignator'][i] == merged_location_top1['RefDesignator'][j]:
                merged_data1_top['Occurence'][i] = merged_location_top1["Occur number"][j]
    merged_data1_top = merged_data1_top.sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
    
    merged_Data_top = merged_data1_top
    merged_Data_top = merged_Data_top.reset_index(drop = True)
else:
    merged_Data_top = pd.DataFrame()

merged_Data_top = merged_Data_top.drop(['Occurence'],axis = 1)

merged_defect_bot = merged_data_bot['Defect'].value_counts()
merged_defect_bot = merged_defect_bot.reset_index()
merged_defect_bot.columns = ["Defect name","Occur number"]
if merged_defect_bot.shape[0] >= 3:
    defect_list_bot = list(merged_defect_bot['Defect name'][0:3])
    
    merged_data1_bot = merged_data_bot[merged_data_bot['Defect'] == defect_list_bot[0]]
    merged_data1_bot = merged_data1_bot.reset_index(drop = True)
    merged_data1_bot['Occurence'] = 1
    merged_location_bot1 = merged_data1_bot['RefDesignator'].value_counts()
    merged_location_bot1 = merged_location_bot1 .reset_index()
    merged_location_bot1.columns = ["RefDesignator","Occur number"]
    for i in range(merged_data1_bot.shape[0]):
        for j in range(merged_location_bot1.shape[0]):
            if merged_data1_bot['RefDesignator'][i] == merged_location_bot1['RefDesignator'][j]:
                merged_data1_bot['Occurence'][i] = merged_location_bot1["Occur number"][j]
    merged_data1_bot = merged_data1_bot.sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
    
    merged_data2_bot = merged_data_bot[merged_data_bot['Defect'] == defect_list_bot[1]]
    merged_data2_bot = merged_data2_bot.reset_index(drop = True)
    merged_data2_bot['Occurence'] = 1
    merged_location_bot2 = merged_data2_bot['RefDesignator'].value_counts()
    merged_location_bot2 = merged_location_bot2 .reset_index()
    merged_location_bot2.columns = ["RefDesignator","Occur number"]
    for i in range(merged_data2_bot.shape[0]):
        for j in range(merged_location_bot2.shape[0]):
            if merged_data2_bot['RefDesignator'][i] == merged_location_bot2['RefDesignator'][j]:
                merged_data2_bot['Occurence'][i] = merged_location_bot2["Occur number"][j]
    merged_data2_bot = merged_data2_bot.sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
    
    merged_data3_bot = merged_data_bot[merged_data_bot['Defect'] == defect_list_bot[2]]
    merged_data3_bot = merged_data3_bot.reset_index(drop = True)
    merged_data3_bot['Occurence'] = 1
    merged_location_bot3 = merged_data3_bot['RefDesignator'].value_counts()
    merged_location_bot3 = merged_location_bot3 .reset_index()
    merged_location_bot3.columns = ["RefDesignator","Occur number"]
    for i in range(merged_data3_bot.shape[0]):
        for j in range(merged_location_bot3.shape[0]):
            if merged_data3_bot['RefDesignator'][i] == merged_location_bot3['RefDesignator'][j]:
                merged_data3_bot['Occurence'][i] = merged_location_bot3["Occur number"][j]
    merged_data3_bot = merged_data3_bot.sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
    merged_Data_bot = pd.concat([merged_data1_bot,merged_data2_bot,merged_data3_bot])
    merged_Data_bot = merged_Data_bot.reset_index(drop = True)
elif merged_defect_bot.shape[0] == 2:
    defect_list_bot = list(merged_defect_bot['Defect name'][0:2])
    
    merged_data1_bot = merged_data_bot[merged_data_bot['Defect'] == defect_list_bot[0]]
    merged_data1_bot = merged_data1_bot.reset_index(drop = True)
    merged_data1_bot['Occurence'] = 1
    merged_location_bot1 = merged_data1_bot['RefDesignator'].value_counts()
    merged_location_bot1 = merged_location_bot1 .reset_index()
    merged_location_bot1.columns = ["RefDesignator","Occur number"]
    for i in range(merged_data1_bot.shape[0]):
        for j in range(merged_location_bot1.shape[0]):
            if merged_data1_bot['RefDesignator'][i] == merged_location_bot1['RefDesignator'][j]:
                merged_data1_bot['Occurence'][i] = merged_location_bot1["Occur number"][j]
    merged_data1_bot = merged_data1_bot.sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
    
    merged_data2_bot = merged_data_bot[merged_data_bot['Defect'] == defect_list_bot[1]]
    merged_data2_bot = merged_data2_bot.reset_index(drop = True)
    merged_data2_bot['Occurence'] = 1
    merged_location_bot2 = merged_data2_bot['RefDesignator'].value_counts()
    merged_location_bot2 = merged_location_bot2 .reset_index()
    merged_location_bot2.columns = ["RefDesignator","Occur number"]
    for i in range(merged_data2_bot.shape[0]):
        for j in range(merged_location_bot2.shape[0]):
            if merged_data2_bot['RefDesignator'][i] == merged_location_bot2['RefDesignator'][j]:
                merged_data2_bot['Occurence'][i] = merged_location_bot2["Occur number"][j]
    merged_data2_bot = merged_data2_bot.sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
    
    merged_Data_bot = pd.concat([merged_data1_bot,merged_data2_bot])
    merged_Data_bot = merged_Data_bot.reset_index(drop = True)
elif merged_defect_bot.shape[0] == 1:
    defect_list_bot = merged_defect_bot['Defect name'][0]
    
    merged_data1_bot = merged_data_bot[merged_data_bot['Defect'] == defect_list_bot]
    merged_data1_bot = merged_data1_bot.reset_index(drop = True)
    merged_data1_bot['Occurence'] = 1
    merged_location_bot1 = merged_data1_bot['RefDesignator'].value_counts()
    merged_location_bot1 = merged_location_bot1 .reset_index()
    merged_location_bot1.columns = ["RefDesignator","Occur number"]
    for i in range(merged_data1_bot.shape[0]):
        for j in range(merged_location_bot1.shape[0]):
            if merged_data1_bot['RefDesignator'][i] == merged_location_bot1['RefDesignator'][j]:
                merged_data1_bot['Occurence'][i] = merged_location_bot1["Occur number"][j]
    merged_data1_bot = merged_data1_bot.sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
    
    merged_Data_bot = merged_data1_bot
    merged_Data_bot = merged_Data_bot.reset_index(drop = True)
else:
    merged_Data_bot = pd.DataFrame()
merged_Data_bot = merged_Data_bot.drop(['Occurence'],axis = 1)
##############################################################################
# If replace, end till here.
##############################################################################

import pptx
from pptx import Presentation
from pptx.util import Pt,Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
ppt= Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[0])

img_path= f"./Bot_{date}_Heatmap.png"
left, top, width, height = Cm(0), Cm(0), Cm(20), Cm(20)

pic= slide.shapes.add_picture(img_path, left, top, width, height) 


if merged_Data_bot.shape[0]>0:

    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    
    left, top, width, height = Cm(0), Cm(0), Cm(30), Cm(20)
    # 表格行列数，和大小
    shape = slide.shapes.add_table(merged_Data_bot.shape[0]+1, merged_Data_bot.shape[1], left, top, width, height)
    # 获取table对象
    table = shape.table
    for i in range(merged_Data_bot.shape[1]):
        table.cell(0, i).text = useful_column[i]
      
    for row in range(1,merged_Data_bot.shape[0] + 1):
            for col in range(0,merged_Data_bot.shape[1]):
                table.cell(row,col).text = str(merged_Data_bot.iloc[row-1,col])      

slide = ppt.slides.add_slide(ppt.slide_layouts[0])

img_path= f"./Top_{date}_Heatmap.png"

left, top, width, height = Cm(0), Cm(0), Cm(20), Cm(20)

pic= slide.shapes.add_picture(img_path, left, top, width, height) 

if merged_Data_top.shape[0]>0:

    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    
    left, top, width, height = Cm(0), Cm(0), Cm(30), Cm(20)
    # 表格行列数，和大小
    shape = slide.shapes.add_table(merged_Data_top.shape[0]+1, merged_Data_top.shape[1], left, top, width, height)
    # 获取table对象
    table = shape.table
    for i in range(merged_Data_top.shape[1]):
        table.cell(0, i).text = useful_column[i]
      
    for row in range(1,merged_Data_top.shape[0] + 1):
            for col in range(0,merged_Data_top.shape[1]):
                table.cell(row,col).text = str(merged_Data_top.iloc[row-1,col])      
                
ppt.save( f'C:/JZhou/L6/Heatmap/heatmap and record till {date}.pptx')

'''
For Shift & Occurence Table, it's not that smart to name the variables based on the top/bot machine,
dynamic naming will be very helpful here, not only save time on naming variables
based on machine but also sequence. Dynamic naming also realize the similar
function with function module.

If you run the code below, no need to run the code between Shift & Occurence 
Table and pptx module(marked)
'''

station = ['top','bot']
#s = 'top'
for s in station:
    locals()['merged_data_' + s] = pd.merge(merged_data, locals()['component_data_'+s], on='RefDesignator')

# Check for mis-matched components
    for index, row in merged_data.iterrows():
        if row['RefDesignator'] not in locals()['merged_data_' + s].values:
            print('\n' + row['RefDesignator'] + ' does not exist in merged_data')
    
    locals()['merged_data_' + s] = locals()['merged_data_' + s][['Label code','Defect','RefDesignator','RegionID_x','Process completion time','Shift']]
    locals()['merged_data_' + s].columns = ['Label code','Defect','RefDesignator','Location','Process completion time','Shift']
    locals()['merged_data_' + s] = locals()['merged_data_' + s].drop_duplicates()

    locals()['merged_defect_' + s] = locals()['merged_data_' + s]['Defect'].value_counts()
    locals()['merged_defect_' + s] = locals()['merged_defect_' + s].reset_index()
    locals()['merged_defect_' + s].columns = ["Defect name","Occur number"]
    
    locals()['merged_Data_' + s] = pd.DataFrame()

    if locals()['merged_defect_' + s].shape[0] != 0:

        for i in range(locals()['merged_defect_' + s].shape[0]):

            locals()['merged_data' + str(i) + '_' + s] = locals()['merged_data_' + s][locals()['merged_data_' + s]['Defect'] == locals()['merged_defect_' + s]['Defect name'][i]]
            locals()['merged_data' + str(i) + '_' + s] = locals()['merged_data' + str(i) + '_' + s].reset_index(drop = True)
            locals()['merged_data' + str(i) + '_' + s]['Occurence'] = 1
    
            locals()['merged_location_' + s + str(i)] = locals()['merged_data' + str(i) + '_' + s]['RefDesignator'].value_counts()
            locals()['merged_location_' + s + str(i)] = locals()['merged_location_' + s + str(i)] .reset_index()
            locals()['merged_location_' + s + str(i)].columns = ["RefDesignator","Occur number"]
            for k in range(locals()['merged_data' + str(i) + '_' + s].shape[0]):
                for j in range(locals()['merged_location_' + s + str(i)].shape[0]):
                    if locals()['merged_data' + str(i) + '_' + s]['RefDesignator'][k] == locals()['merged_location_' + s + str(i)]['RefDesignator'][j]:
                        locals()['merged_data' + str(i) + '_' + s]['Occurence'][k] = locals()['merged_location_' + s + str(i)]["Occur number"][j]
            locals()['merged_data' + str(i) + '_' + s] = locals()['merged_data' + str(i) + '_' + s].sort_values(by = ['Occurence','RefDesignator'],ascending = (False,True))
            
            locals()['merged_Data_' + s] = pd.concat([locals()['merged_Data_' + s],locals()['merged_data' + str(i) + '_' + s]])
            locals()['merged_Data_' + s] = locals()['merged_Data_' + s].reset_index(drop = True)
            if i == 2:
                break
    locals()['merged_Data_' + s] = locals()['merged_Data_' + s].drop(['Occurence'],axis = 1)


